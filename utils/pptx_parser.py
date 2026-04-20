"""
PPTX Parser — extracts text content from each slide.
"""

import base64
import io
from pptx import Presentation


def parse_pptx(contents_b64: str) -> list[dict]:
    """
    Decode a base64-encoded PPTX file (from Dash Upload),
    extract text from every slide, and return a list of dicts:
        [{"slide": 1, "text": "All text on slide 1 ..."}, ...]
    """
    # Strip the data-URI prefix if present
    if "," in contents_b64:
        contents_b64 = contents_b64.split(",", 1)[1]

    raw = base64.b64decode(contents_b64)
    prs = Presentation(io.BytesIO(raw))

    slides = []
    for idx, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        parts.append(text)
            # Also capture table text
            if shape.has_table:
                for row in shape.table.rows:
                    row_texts = [
                        cell.text.strip() for cell in row.cells if cell.text.strip()
                    ]
                    if row_texts:
                        parts.append(" | ".join(row_texts))

        slides.append({"slide": idx, "text": "\n".join(parts)})

    return slides


def format_slides_for_prompt(slides: list[dict]) -> str:
    """Format parsed slides into a readable string for LLM context."""
    if not slides:
        return "(No slides loaded)"
    shapes = []
    lines = []
    for slide_shapes in slides:
        shapes.extend(slides[slide_shapes])
        lines.append(f"--- Slide {slide_shapes} ---")
        for s in shapes:
            lines.append(s["full_text"] if s["full_text"] else "")
            lines.append("")

    return "\n".join(lines)
