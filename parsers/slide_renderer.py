"""
CiteMind — HTML/CSS slide renderer.

Converts a single PPTX slide to a self-contained HTML string using python-pptx.
Each shape is rendered as an absolutely-positioned div with background/border/text
derived directly from the EMU-level data.  No LibreOffice / soffice required.
"""

from __future__ import annotations

import io
import json
import logging
import re
from typing import Any, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

log = logging.getLogger(__name__)


# ─────────────────────────────────────────────────────────────────────────────
# Internal helpers
# ─────────────────────────────────────────────────────────────────────────────


def _emu_to_pct(val: int, total: int) -> float:
    """Convert an EMU value to percentage of the slide dimension."""
    return (val / total * 100) if total else 0


def _rgb_to_hex(rgb: Optional[RGBColor]) -> Optional[str]:
    if rgb is None:
        return None
    return "#{:02X}{:02X}{:02X}".format(rgb[0], rgb[1], rgb[2])


def _shape_fill_css(shape: Any) -> str:
    """Return a CSS background-color string for a shape fill, or ''."""
    try:
        fill = shape.fill
        if fill.type is None:
            return ""
        # Solid fill
        from pptx.enum.dml import MSO_THEME_COLOR

        fg = fill.fore_color
        try:
            rgb = fg.rgb
            return "background-color: {};".format(_rgb_to_hex(rgb))
        except Exception:
            return ""
    except Exception:
        return ""


def _font_size_pt(run: Any) -> Optional[float]:
    """Return font size in points, or None."""
    try:
        if run.font.size:
            # Scale down slightly (~85%) since web rendering of pt 
            # often appears larger than native PPTX on desktop windows
            return run.font.size.pt * 0.85
    except Exception:
        pass
    return None


def _run_to_span(run: Any) -> str:
    """Convert a pptx Run to an HTML <span>."""
    styles = []
    try:
        if run.font.bold:
            styles.append("font-weight:bold")
        if run.font.italic:
            styles.append("font-style:italic")
        if run.font.underline:
            styles.append("text-decoration:underline")
        pt = _font_size_pt(run)
        if pt:
            styles.append("font-size:{:.2f}pt".format(pt))
        try:
            rgb = run.font.color.rgb
            styles.append("color:{}".format(_rgb_to_hex(rgb)))
        except Exception:
            pass
    except Exception:
        pass

    style_str = ";".join(styles)
    text = _escape_html(run.text)
    return '<span style="{}">{}</span>'.format(style_str, text)


def _para_align_css(para: Any) -> str:
    try:
        align = para.alignment
        if align == PP_ALIGN.CENTER:
            return "text-align:center;"
        if align == PP_ALIGN.RIGHT:
            return "text-align:right;"
        if align == PP_ALIGN.JUSTIFY:
            return "text-align:justify;"
    except Exception:
        pass
    return "text-align:left;"


def _escape_html(text: str) -> str:
    return (
        text.replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
        .replace('"', "&quot;")
    )


def _text_frame_html(shape: Any) -> str:
    """Render a text frame as HTML paragraphs."""
    if not shape.has_text_frame:
        return ""
    parts = []
    for para in shape.text_frame.paragraphs:
        align_css = _para_align_css(para)
        spans = "".join(_run_to_span(r) for r in para.runs if r.text)
        # Paragraph spacing / indent
        parts.append(
            '<p style="margin:0;padding:0;{}">{}</p>'.format(
                align_css, spans or "&nbsp;"
            )
        )
    return "".join(parts)


def _table_html(shape: Any) -> str:
    """Render a pptx table as an HTML table."""
    rows_html = []
    for row in shape.table.rows:
        cells_html = []
        for cell in row.cells:
            text = _escape_html(cell.text.strip())
            cells_html.append(
                "<td style='padding:4px;border:1px solid #ccc'>{}</td>".format(text)
            )
        rows_html.append("<tr>{}</tr>".format("".join(cells_html)))
    return (
        '<table style="border-collapse:collapse;width:100%;height:100%;'
        'font-size:11pt;table-layout:fixed">{}</table>'.format("".join(rows_html))
    )


# ─────────────────────────────────────────────────────────────────────────────
# Public API
# ─────────────────────────────────────────────────────────────────────────────


def render_slide_to_html(pptx_bytes: bytes, slide_index: int) -> str:
    """
    Parse *pptx_bytes* with python-pptx and return a self-contained HTML
    fragment that renders slide *slide_index* using absolute-% positioning.

    The returned string is suitable for embedding inside a fixed-size container
    with ``position:relative`` and a 16:9 (or actual slide) aspect ratio.
    """
    try:
        prs = Presentation(io.BytesIO(pptx_bytes))
        slide_w = int(prs.slide_width)
        slide_h = int(prs.slide_height)

        slides = prs.slides
        if slide_index < 0 or slide_index >= len(slides):
            return _error_html("Slide index {} out of range".format(slide_index))

        slide = slides[slide_index]
        shapes_html = []

        # ── slide background ──────────────────────────────────────────────────
        bg_css = _slide_background_css(slide)

        for shape in slide.shapes:
            try:
                left_pct = _emu_to_pct(shape.left or 0, slide_w)
                top_pct = _emu_to_pct(shape.top or 0, slide_h)
                width_pct = _emu_to_pct(shape.width or 0, slide_w)
                height_pct = _emu_to_pct(shape.height or 0, slide_h)

                fill_css = _shape_fill_css(shape)

                # Inner content
                if hasattr(shape, "has_table") and shape.has_table:
                    inner = _table_html(shape)
                elif shape.has_text_frame:
                    inner = _text_frame_html(shape)
                else:
                    inner = ""

                div = (
                    '<div style="'
                    "position:absolute;"
                    "left:{left:.4f}%;"
                    "top:{top:.4f}%;"
                    "width:{width:.4f}%;"
                    "height:{height:.4f}%;"
                    "overflow:hidden;"
                    "box-sizing:border-box;"
                    "{fill}"
                    '">'
                    "{inner}"
                    "</div>"
                ).format(
                    left=left_pct,
                    top=top_pct,
                    width=width_pct,
                    height=height_pct,
                    fill=fill_css,
                    inner=inner,
                )
                shapes_html.append(div)
            except Exception as shape_err:
                log.debug(
                    "Skipped shape %r: %s", getattr(shape, "name", "?"), shape_err
                )

        return (
            '<div style="width:100%;aspect-ratio:16/9;position:relative;">'
            '<div style="'
            "position:absolute;"
            "inset:0;"
            "overflow:hidden;"
            "font-family:Calibri,Arial,sans-serif;"
            "font-size:14pt;"
            "{bg}"
            '">'
            "{shapes}"
            "</div>"
            "</div>"
        ).format(bg=bg_css, shapes="".join(shapes_html))

    except Exception as e:
        log.error("render_slide_to_html failed: %s", e)
        return _error_html(str(e))


def _slide_background_css(slide: Any) -> str:
    """Return CSS for slide background (solid color or transparent)."""
    try:
        bg = slide.background
        fill = bg.fill
        if fill.type is not None:
            fg = fill.fore_color
            rgb = fg.rgb
            return "background-color:{};".format(_rgb_to_hex(rgb))
    except Exception:
        pass
    return "background-color:#FFFFFF;"


def _error_html(msg: str) -> str:
    return (
        '<div style="display:flex;align-items:center;justify-content:center;'
        'width:100%;height:100%;color:#f87171;font-size:13pt;">'
        "⚠ Could not render slide: {}</div>"
    ).format(_escape_html(msg))
