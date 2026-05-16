"""
CiteMind — PowerPoint parser with shape-level extraction.
Extracts every shape with normalised coordinates and run-level text.
"""
from __future__ import annotations

import base64
import hashlib
import io
import json
import re
from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.util import Emu


def serialize_slide(slide: Any, slide_w: int, slide_h: int) -> List[dict]:
    """
    Return list of shape dicts with normalised coordinates (0.0-1.0).

    Each dict has:
        pptx_shape_id, shape_name, shape_type, x_pct, y_pct, w_pct, h_pct,
        full_text, runs_json, z_order, has_table, has_chart
    """
    shapes = []  # type: List[dict]
    for z_order, shape in enumerate(slide.shapes):
        entry = {
            "pptx_shape_id": shape.shape_id,
            "shape_name": shape.name,
            "shape_type": str(shape.shape_type).split(".")[-1],
            "x_pct": shape.left / slide_w if slide_w else 0,
            "y_pct": shape.top / slide_h if slide_h else 0,
            "w_pct": shape.width / slide_w if slide_w else 0,
            "h_pct": shape.height / slide_h if slide_h else 0,
            "z_order": z_order,
            "full_text": "",
            "runs": [],
            "has_table": False,
            "has_chart": False,
        }

        if shape.has_text_frame:
            run_idx = 0
            all_text = []  # type: List[str]
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    font_size = None
                    if run.font.size is not None:
                        font_size = int(run.font.size)
                    entry["runs"].append({
                        "index": run_idx,
                        "text": run.text,
                        "bold": bool(run.font.bold),
                        "size": font_size,
                        "is_numeric": bool(re.search(r'[\d,.%₹$€£¥]', run.text)),
                    })
                    all_text.append(run.text)
                    run_idx += 1
            entry["full_text"] = "".join(all_text)

        if hasattr(shape, "has_table") and shape.has_table:
            entry["has_table"] = True
            # Extract table text as runs too
            run_idx = len(entry["runs"])
            table_texts = []  # type: List[str]
            for row in shape.table.rows:
                for cell in row.cells:
                    text = cell.text.strip()
                    if text:
                        entry["runs"].append({
                            "index": run_idx,
                            "text": text,
                            "bold": False,
                            "size": None,
                            "is_numeric": bool(re.search(r'[\d,.%₹$€£¥]', text)),
                        })
                        table_texts.append(text)
                        run_idx += 1
            if table_texts:
                if entry["full_text"]:
                    entry["full_text"] += " " + " ".join(table_texts)
                else:
                    entry["full_text"] = " ".join(table_texts)

        if hasattr(shape, "has_chart") and shape.has_chart:
            entry["has_chart"] = True

        entry["runs_json"] = json.dumps(entry.pop("runs"))
        shapes.append(entry)
    return shapes


def parse_pptx_file(contents_b64: str) -> dict:
    """
    Parse a base64-encoded PPTX (from Dash Upload).

    Returns dict:
        {
            "sha256": str,
            "slide_width_emu": int,
            "slide_height_emu": int,
            "slide_count": int,
            "slides": [
                {
                    "slide_index": int,
                    "slide_number": int,
                    "title": str | None,
                    "shapes": [shape_dict, ...],
                    "has_table": bool,
                    "has_chart": bool,
                    "shape_count": int,
                }
            ]
        }
    """
    # Strip data-URI prefix
    if "," in contents_b64:
        contents_b64 = contents_b64.split(",", 1)[1]

    raw = base64.b64decode(contents_b64)
    sha = hashlib.sha256(raw).hexdigest()
    prs = Presentation(io.BytesIO(raw))

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    slides = []  # type: List[dict]
    for idx, slide in enumerate(prs.slides):
        shapes = serialize_slide(slide, slide_w, slide_h)

        # Try to find slide title
        title = None
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            title = slide.shapes.title.text_frame.text.strip() or None

        has_table = any(s.get("has_table") for s in shapes)
        has_chart = any(s.get("has_chart") for s in shapes)

        slides.append({
            "slide_index": idx,
            "slide_number": idx + 1,
            "title": title,
            "shapes": shapes,
            "has_table": has_table,
            "has_chart": has_chart,
            "shape_count": len(shapes),
        })

    return {
        "sha256": sha,
        "slide_width_emu": int(slide_w),
        "slide_height_emu": int(slide_h),
        "slide_count": len(slides),
        "slides": slides,
    }


def get_raw_bytes(contents_b64: str) -> bytes:
    """Decode base64 content from Dash Upload to raw bytes."""
    if "," in contents_b64:
        contents_b64 = contents_b64.split(",", 1)[1]
    return base64.b64decode(contents_b64)


def format_slides_for_prompt(slides_data: List[dict]) -> str:
    """Format parsed slides into a readable string for LLM context."""
    if not slides_data:
        return "(No slides loaded)"

    lines = []  # type: List[str]
    for s in slides_data:
        title = s.get("title") or "(untitled)"
        lines.append("--- Slide {} ({}) ---".format(s["slide_number"], title))
        shapes = s.get("shapes", [])
        for shape in shapes:
            text = shape.get("full_text", "")
            if text and text.strip():
                lines.append("  [{}] {}".format(shape["shape_name"], text.strip()))
        lines.append("")
    return "\n".join(lines)
